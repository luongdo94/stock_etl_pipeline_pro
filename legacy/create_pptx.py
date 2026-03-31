"""
create_pptx.py — Generate the Elite Pro Stock Analytics Presentation.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import pandas as pd
import duckdb

ROOT = os.path.dirname(os.path.abspath(__file__))

# ── Color palette ──────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0A, 0x0A, 0x0A)   # Black
BG_CARD    = RGBColor(0x1E, 0x1E, 0x1E)   # Dark Gray
ACCENT     = RGBColor(0x00, 0xD4, 0xFF)   # Cyber Blue
ACCENT2    = RGBColor(0x00, 0xFF, 0x88)   # Neon Green
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0xB0, 0xB0, 0xB0)
YELLOW     = RGBColor(0xFF, 0xCC, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

def add_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

# ── SLIDE 1: Title ─────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
add_bg(s1)
add_rect(s1, 0, 3.4, 13.33, 0.1, ACCENT)
add_text(s1, "🚀", 6.1, 1.2, 1, 1, size=60, align=PP_ALIGN.CENTER)
add_text(s1, "Elite Pro Stock Analytics", 1.2, 2.2, 11, 1, size=44, bold=True, align=PP_ALIGN.CENTER)
add_text(s1, "End-to-End Modern Data Stack for Financial Markets", 1.2, 3.8, 11, 0.5, size=22, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s1, "Python · DuckDB · Airflow · Streamlit · GitHub Actions", 1.2, 4.5, 11, 0.5, size=16, color=GRAY, align=PP_ALIGN.CENTER)

# ── SLIDE 2: Project Overview ──────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
add_bg(s2)
add_text(s2, "🎯 Project Overview", 0.5, 0.3, 10, 0.8, size=32, bold=True, color=ACCENT)
items = [
    ("💎 Universe", "Tracking 20 global Tech/Healthcare giants (US, EU, CN)."),
    ("📊 Fundamental", "4-Year Revenue & EPS History with YoY growth analysis."),
    ("🧠 AI Engine", "Unified 9-factor scoring model (Valuation, Growth, Techs)."),
    ("🏦 Warehouse", "Medallion Architecture (Bronze/Silver/Gold) on DuckDB."),
    ("⏲️ Automation", "Fully automated daily updates via Airflow & GitHub Actions."),
    ("📈 Reporting", "Rich HTML Email reports & Interactive Streamlit Dashboard."),
]
for i, (ttl, desc) in enumerate(items):
    y = 1.5 + i * 0.95
    add_text(s2, f"● {ttl}:", 0.6, y, 3, 0.4, size=18, bold=True, color=ACCENT2)
    add_text(s2, desc, 3.5, y, 9, 0.4, size=18, color=WHITE)

# ── SLIDE 3: Technical Stack ───────────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
add_bg(s3)
add_text(s3, "🏗️ Technical Stack (The 'Pro' Layer)", 0.5, 0.3, 10, 0.8, size=32, bold=True, color=ACCENT)
techs = [
    ("Python/yfinance", "Robust data extraction with automatic NaN filtering."),
    ("DuckDB", "High-performance analytics engine (Columnar storage)."),
    ("Airflow/Docker", "Orchestration & reproducibility in any environment."),
    ("Streamlit", "Premium Glassmorphism UI for real-time visualization."),
    ("GitHub Actions", "24/7 Serverless automation and CI/CD sync."),
]
for i, (t, d) in enumerate(techs):
    y = 1.8 + i * 1.0
    add_rect(s3, 0.5, y, 12.3, 0.8, BG_CARD)
    add_text(s3, t, 0.7, y + 0.15, 3.5, 0.5, size=18, bold=True, color=ACCENT2)
    add_text(s3, d, 4.5, y + 0.15, 8, 0.5, size=16, color=WHITE)

# ── SLIDE 4: Live Insights (FETCH FROM DUCKDB) ─────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
add_bg(s4)
add_text(s4, "📊 Current Market Insights (Live Data)", 0.5, 0.3, 10, 0.8, size=32, bold=True, color=YELLOW)

db_path = os.path.join(ROOT, "warehouse", "stock_dw.duckdb")
try:
    conn = duckdb.connect(db_path, read_only=True)
    latest = conn.execute("""
        SELECT ticker, ma_signal, price_close, pct_from_52w_high
        FROM marts.fct_daily_returns 
        WHERE date = (SELECT MAX(date) FROM marts.fct_daily_returns)
    """).df()
    perf = conn.execute("""
        SELECT ticker, ROUND(AVG(daily_return_pct), 2) as avg_ret
        FROM marts.fct_daily_returns 
        WHERE date >= (SELECT MAX(date) - INTERVAL 30 DAY FROM marts.fct_daily_returns)
        GROUP BY ticker ORDER BY avg_ret DESC
    """).df()
    conn.close()

    if not latest.empty:
        top_g = f"{perf.iloc[0]['ticker']} (+{perf.iloc[0]['avg_ret']}% avg daily)"
        bullish = len(latest[latest['ma_signal'] == 'BULLISH'])
        total = len(latest)
        insights = [
            ("🏆 Top Performer", top_g),
            ("📡 Trend Guard", f"{bullish}/{total} stocks in Bullish trend (MA20 > MA50)."),
            ("📏 High Ground", f"{latest.sort_values('pct_from_52w_high', ascending=False).iloc[0]['ticker']} is closest to 52w high."),
            ("🛡️ Deep Value", f"Only {len(latest[latest['pct_from_52w_high'] < -20])} stocks corrected >20%."),
        ]
    else:
        insights = [("⚠️ Status", "No daily data found in warehouse.")]
except Exception as e:
    insights = [("❌ Data Error", str(e))]

for i, (t, d) in enumerate(insights):
    y = 1.8 + i * 1.3
    add_rect(s4, 1, y, 11, 1.1, BG_CARD)
    add_text(s4, t, 1.2, y + 0.15, 4, 0.5, size=22, bold=True, color=ACCENT2)
    add_text(s4, d, 1.2, y + 0.6, 10, 0.5, size=18, color=WHITE)

# ── SLIDE 5: Conclusion ────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
add_bg(s5)
add_text(s5, "Thank You! / Cảm ơn bạn!", 1.2, 2.5, 11, 1, size=44, bold=True, align=PP_ALIGN.CENTER)
add_text(s5, "GitHub: github.com/luongdo94/stock_etl_pipeline_pro", 1.2, 4.0, 11, 0.5, size=18, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s5, "Questions & Discussion", 1.2, 5.0, 11, 0.5, size=16, color=GRAY, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
OUTPUT = os.path.join(ROOT, "Elite_Pro_Dashboard_Presentation.pptx")
prs.save(OUTPUT)
print(f"✅ PowerPoint saved: {OUTPUT}")

# Open on Mac
import subprocess
try:
    subprocess.Popen(["open", OUTPUT])
    print("🚀 Opening presentation...")
except:
    pass
