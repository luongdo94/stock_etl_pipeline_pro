"""
create_pptx.py — Generate a German interview presentation PowerPoint.
Usage: python c:\\etl_pipeline\\create_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Color palette ──────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0D, 0x1B, 0x2A)   # Deep navy
BG_CARD    = RGBColor(0x1A, 0x2A, 0x3A)   # Card background
ACCENT     = RGBColor(0x00, 0xC8, 0xFF)   # Cyan accent
ACCENT2    = RGBColor(0x00, 0xFF, 0xA3)   # Green accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x8A, 0x9B, 0xAD)
YELLOW     = RGBColor(0xFF, 0xD6, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # Blank layout


def add_bg(slide, color=BG_DARK):
    """Fill slide background with solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txb


def add_multiline(slide, lines, l, t, w, h, size=14, color=WHITE, spacing=0.35):
    """Add multiple lines as separate text boxes stacked vertically."""
    for i, (bullet, text, clr) in enumerate(lines):
        y = t + i * spacing
        add_text(slide, f"{bullet}  {text}", l, y, w, 0.35, size=size, color=clr)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1: Title
# ─────────────────────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
add_bg(s1)

# Decorative accent bar
add_rect(s1, 0, 0, 0.5, 7.5, ACCENT)
add_rect(s1, 0.5, 3.4, 12.83, 0.08, ACCENT)

add_text(s1, "📊", 1.2, 1.0, 2, 1.2, size=60, align=PP_ALIGN.LEFT)
add_text(s1, "Stock Market ETL Pipeline", 1.2, 2.0, 11, 1.2,
         size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s1, "End-to-End Data Engineering Projekt", 1.2, 3.1, 11, 0.6,
         size=22, color=ACCENT, align=PP_ALIGN.LEFT)
add_text(s1, "Python  ·  DuckDB  ·  Apache Airflow  ·  Docker  ·  Plotly",
         1.2, 3.7, 11, 0.5, size=16, color=GRAY, align=PP_ALIGN.LEFT)
add_text(s1, "Präsentation auf Deutsch", 1.2, 6.5, 6, 0.5,
         size=13, color=GRAY, align=PP_ALIGN.LEFT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2: Projektüberblick
# ─────────────────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
add_bg(s2)
add_rect(s2, 0, 0, 13.33, 1.1, BG_CARD)
add_rect(s2, 0, 1.08, 13.33, 0.05, ACCENT)
add_text(s2, "🎯  Projektüberblick", 0.4, 0.2, 12, 0.7,
         size=28, bold=True, color=WHITE)

items = [
    ("▶", "Automatische tägliche Datenerfassung für 8 internationale Tech-Aktien"),
    ("▶", "Datenquellen: Yahoo Finance API (kostenlos, real-time)"),
    ("▶", "Lokales Data Warehouse mit DuckDB nach Medallion Architecture"),
    ("▶", "4 Datenschichten: Raw → Staging → Intermediate → Marts"),
    ("▶", "Automatisierung mit Apache Airflow (täglich 18:00 nach NYSE-Schluss)"),
    ("▶", "Interaktives Dashboard mit 5 Analysegrafiken in Plotly"),
    ("▶", "Deployment via Docker – vollständig reproduzierbar"),
]
for i, (b, t) in enumerate(items):
    add_text(s2, f"{b}  {t}", 0.5, 1.4 + i * 0.77, 12.3, 0.6, size=15, color=WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3: Daten & Abdeckung
# ─────────────────────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
add_bg(s3)
add_rect(s3, 0, 0, 13.33, 1.1, BG_CARD)
add_rect(s3, 0, 1.08, 13.33, 0.05, ACCENT)
add_text(s3, "📈  Daten & Abdeckung", 0.4, 0.2, 12, 0.7, size=28, bold=True, color=WHITE)

tickers = [
    ("NVDA",  "NVIDIA",    "Halbleiter",     "USA",  "🇺🇸"),
    ("MSFT",  "Microsoft", "Cloud/Software", "USA",  "🇺🇸"),
    ("AAPL",  "Apple",     "Consumer Tech",  "USA",  "🇺🇸"),
    ("GOOGL", "Alphabet",  "Cloud/Software", "USA",  "🇺🇸"),
    ("AMZN",  "Amazon",    "Cloud/E-Com",    "USA",  "🇺🇸"),
    ("SAP",   "SAP SE",    "Enterprise SW",  "EU",   "🇩🇪"),
    ("ASML",  "ASML",      "Halbleiter",     "EU",   "🇳🇱"),
    ("BABA",  "Alibaba",   "E-Commerce",     "CN",   "🇨🇳"),
]
cols = [0.5, 2.0, 4.2, 6.5, 8.5]
headers = ["Ticker", "Unternehmen", "Sektor", "Region"]
for j, h in enumerate(headers):
    add_text(s3, h, cols[j], 1.3, 2.2, 0.4, size=13, bold=True, color=ACCENT)

for i, (ticker, name, sector, region, flag) in enumerate(tickers):
    y = 1.8 + i * 0.62
    if i % 2 == 0:
        add_rect(s3, 0.3, y - 0.05, 12.7, 0.55, BG_CARD)
    add_text(s3, ticker, cols[0], y, 1.4, 0.5, size=14, bold=True, color=ACCENT2)
    add_text(s3, name,   cols[1], y, 2.0, 0.5, size=14, color=WHITE)
    add_text(s3, sector, cols[2], y, 2.2, 0.5, size=14, color=WHITE)
    add_text(s3, f"{flag} {region}", cols[3], y, 1.5, 0.5, size=14, color=WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4: Technische Architektur
# ─────────────────────────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
add_bg(s4)
add_rect(s4, 0, 0, 13.33, 1.1, BG_CARD)
add_rect(s4, 0, 1.08, 13.33, 0.05, ACCENT)
add_text(s4, "🏗️  Technische Architektur – Datenfluss", 0.4, 0.2, 12, 0.7,
         size=28, bold=True, color=WHITE)

steps = [
    ("Yahoo Finance\nAPI", ACCENT),
    ("EXTRACT\nextract.py", RGBColor(0x1E,0x88,0xE5)),
    ("LOAD\nload.py", RGBColor(0x43,0xA0,0x47)),
    ("TRANSFORM\ntransform.py", RGBColor(0xFB,0x8C,0x00)),
    ("MARTS\nDuckDB", RGBColor(0x8E,0x24,0xAA)),
    ("DASHBOARD\nPlotly", RGBColor(0xE5,0x39,0x35)),
    ("AIRFLOW\nDocker", ACCENT),
]
box_w, box_h = 1.55, 1.0
start_x = 0.3
y_pos = 2.5
for i, (label, color) in enumerate(steps):
    x = start_x + i * (box_w + 0.22)
    add_rect(s4, x, y_pos, box_w, box_h, color)
    add_text(s4, label, x, y_pos + 0.15, box_w, box_h - 0.1,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text(s4, "→", x + box_w, y_pos + 0.28, 0.25, 0.5,
                 size=22, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

add_text(s4, "Medallion Architecture: raw → staging → intermediate → marts",
         1.0, 4.0, 11.3, 0.5, size=15, color=ACCENT, align=PP_ALIGN.CENTER)

layers = [
    ("🟤 RAW",          "Rohdaten – unveränderter Originalstand"),
    ("🔵 STAGING",      "Datenbereinigung & Validierung"),
    ("🟠 INTERMEDIATE", "Technische Indikatoren: MA7/20/50, tägliche Rendite, Volatilität"),
    ("🟢 MARTS",        "Endprodukt-Tabellen für Analysen & BI-Tools"),
]
for i, (layer, desc) in enumerate(layers):
    add_text(s4, f"{layer}: {desc}", 0.5, 4.7 + i * 0.55, 12.3, 0.5,
             size=13, color=WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5: Technologie-Stack
# ─────────────────────────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
add_bg(s5)
add_rect(s5, 0, 0, 13.33, 1.1, BG_CARD)
add_rect(s5, 0, 1.08, 13.33, 0.05, ACCENT)
add_text(s5, "⚙️  Technologie-Stack & Entscheidungsbegründung", 0.4, 0.2,
         12, 0.7, size=28, bold=True, color=WHITE)

techs = [
    ("Python 3.11",      "Industriestandard für Data Engineering. Breites Ökosystem."),
    ("DuckDB",           "OLAP-optimierte Datenbank – columnar, schnell, serverless."),
    ("Apache Airflow",   "De-facto-Standard für Workflow-Orchestrierung in der Industrie."),
    ("Docker",           "Reproduzierbare Umgebung – läuft auf jedem System identisch."),
    ("Plotly",           "Interaktive Visualisierungen ohne Frontend-Kenntnisse."),
    ("yfinance",         "Kostenloser Zugang zu realen Börsendaten von Yahoo Finance."),
]
for i, (tech, desc) in enumerate(techs):
    col = 0 if i < 3 else 6.8
    row = i % 3
    y = 1.5 + row * 1.6
    add_rect(s5, col + 0.3, y, 6.0, 1.3, BG_CARD)
    add_text(s5, tech, col + 0.6, y + 0.1, 5.5, 0.5, size=16, bold=True, color=ACCENT2)
    add_text(s5, desc, col + 0.6, y + 0.55, 5.5, 0.65, size=13, color=WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6: Code-Qualität & Best Practices
# ─────────────────────────────────────────────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
add_bg(s6)
add_rect(s6, 0, 0, 13.33, 1.1, BG_CARD)
add_rect(s6, 0, 1.08, 13.33, 0.05, ACCENT)
add_text(s6, "✅  Code-Qualität & Best Practices", 0.4, 0.2, 12, 0.7,
         size=28, bold=True, color=WHITE)

practices = [
    ("🧩", "Modularität",            "Jede Schicht ist eine eigene Datei (SRP-Prinzip)"),
    ("🔄", "Upsert-Logik",           "Verhindert doppelte Einträge bei täglichem Lauf"),
    ("🛡️", "Data Quality Checks",   "4 automatische Tests nach jeder Transformation"),
    ("📝", "Strukturiertes Logging", "Zeitstempel, Row Counts, Step-Timing in jeder Phase"),
    ("⚡", "Performance",            "DuckDB lädt 2.000 Zeilen in < 0,1 Sekunden"),
    ("🔐", "Sichere Konfiguration",  "SMTP-Passwörter als Docker-Umgebungsvariablen"),
    ("🔁", "Idempotenz",             "Pipeline kann mehrfach ohne Fehler ausgeführt werden"),
    ("🐳", "Docker-Ready",           "Einziger Befehl: docker-compose up airflow -d"),
]
for i, (icon, title, desc) in enumerate(practices):
    col = 0.4 if i < 4 else 7.0
    row = i % 4
    y = 1.4 + row * 1.45
    add_rect(s6, col, y, 6.1, 1.2, BG_CARD)
    add_text(s6, f"{icon}  {title}", col + 0.2, y + 0.08, 5.7, 0.45,
             size=15, bold=True, color=ACCENT2)
    add_text(s6, desc, col + 0.2, y + 0.55, 5.7, 0.55, size=13, color=WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7: Live-Demo Ablauf
# ─────────────────────────────────────────────────────────────────────────────
s7 = prs.slides.add_slide(BLANK)
add_bg(s7)
add_rect(s7, 0, 0, 13.33, 1.1, BG_CARD)
add_rect(s7, 0, 1.08, 13.33, 0.05, ACCENT)
add_text(s7, "💻  Live-Demo Ablauf", 0.4, 0.2, 12, 0.7, size=28, bold=True, color=WHITE)

steps7 = [
    ("1", "Pipeline starten",    "python c:\\etl_pipeline\\run.py",
     "4 Phasen: Extract → Validate → Load → Transform"),
    ("2", "Daten abfragen",      "duckdb -readonly ... + SQL-Abfragen",
     "Top-Performer der letzten 30 Tage anzeigen"),
    ("3", "Dashboard öffnen",    "python c:\\etl_pipeline\\dashboard.py",
     "5 interaktive Charts im Browser präsentieren"),
    ("4", "Airflow UI zeigen",   "http://localhost:8080",
     "DAG-Graph, Task-Logs, Ausführungshistorie"),
]
for i, (num, title, cmd, desc) in enumerate(steps7):
    y = 1.4 + i * 1.35
    add_rect(s7, 0.4, y, 0.65, 0.65, ACCENT)
    add_text(s7, num, 0.4, y + 0.05, 0.65, 0.6, size=22, bold=True,
             color=BG_DARK, align=PP_ALIGN.CENTER)
    add_text(s7, title, 1.3, y + 0.02, 5.0, 0.45, size=16, bold=True, color=WHITE)
    add_rect(s7, 1.3, y + 0.5, 11.0, 0.38, RGBColor(0x0A, 0x0A, 0x1A))
    add_text(s7, f"  {cmd}", 1.3, y + 0.5, 6.5, 0.38, size=12, color=ACCENT2)
    add_text(s7, f"→ {desc}", 8.0, y + 0.5, 4.6, 0.38, size=12, color=GRAY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8: Ergebnisse & Insights (DYNAMIC FROM DUCKDB)
# ─────────────────────────────────────────────────────────────────────────────
s8 = prs.slides.add_slide(BLANK)
add_bg(s8)
add_rect(s8, 0, 0, 13.33, 1.1, BG_CARD)
add_rect(s8, 0, 1.08, 13.33, 0.05, ACCENT)
add_text(s8, "📊  Aktuelle Marktanalyse – Key Insights", 0.4, 0.2, 12, 0.7,
         size=28, bold=True, color=WHITE)

# Fetch actual live data from DuckDB for presentation!
import duckdb
db_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "warehouse", "stock_dw.duckdb")
try:
    conn = duckdb.connect(db_path, read_only=True)
    monthly = conn.execute("""
        SELECT f.ticker, d.company, ROUND(SUM(f.daily_return_pct), 2) AS monthly_pct
        FROM marts.fct_daily_returns f LEFT JOIN marts.dim_companies d USING (ticker)
        WHERE f.date >= (SELECT MAX(date) - INTERVAL 30 DAY FROM marts.fct_daily_returns)
        GROUP BY f.ticker, d.company ORDER BY monthly_pct DESC
    """).df()
    
    latest = conn.execute("""
        SELECT ticker, ma_signal, ROUND(pct_from_52w_high, 2) AS pct_from_52w_high 
        FROM marts.fct_daily_returns WHERE date = (SELECT MAX(date) FROM marts.fct_daily_returns)
    """).df()
    conn.close()

    total_t     = len(monthly)
    neg_count   = len(monthly[monthly['monthly_pct'] < 0])
    best_p      = f"{monthly.iloc[0]['ticker']} +{monthly.iloc[0]['monthly_pct']}%" if total_t > 0 else "N/A"
    worst_p1    = f"{monthly.iloc[-1]['ticker']} {monthly.iloc[-1]['monthly_pct']}%" if total_t > 0 else "N/A"
    worst_p2    = f" & {monthly.iloc[-2]['ticker']} {monthly.iloc[-2]['monthly_pct']}%" if total_t > 1 else ""
    bearish     = len(latest[latest['ma_signal'] == 'BEARISH'])
    bullish_arr = latest[latest['ma_signal'] == 'BULLISH']['ticker'].tolist()
    bullish_str = ", ".join(bullish_arr) if bullish_arr else "Kein Titel"
    worst_52w   = latest.sort_values("pct_from_52w_high").iloc[0]

    insights = [
        ("📉", "Marktkorrektur",     f"{neg_count} von {total_t} Titeln mit negativer 30-Tage-Rendite"),
        ("🏅", "Bester Performer",   f"{best_p} in den letzten 30 Tagen"),
        ("⚠️", "Stärkste Korrektur",f"{worst_p1}{worst_p2} in den letzten 30 Tagen"),
        ("📡", "MA-Signal",          f"{bearish}/{total_t} Titel BEARISH (MA20 < MA50) – Abwärtstrend"),
        ("💎", "Technisch stärkstes",f"{bullish_str} zeigen aktuell ein BULLISH-Signal"),
        ("📏", "52-Wochen-Tief",     f"{worst_52w['ticker']} liegt {abs(worst_52w['pct_from_52w_high'])}% unter dem Jahreshoch"),
    ]
except Exception as e:
    insights = [("❌", "Datenfehler", f"DuckDB konnte nicht gelesen werden: {e}")]

for i, (icon, title, desc) in enumerate(insights):
    if i >= 6: break
    col = 0.4 if i < 3 else 6.9
    row = i % 3
    y = 1.4 + row * 1.8
    add_rect(s8, col, y, 5.9, 1.5, BG_CARD)
    add_text(s8, f"{icon}  {title}", col + 0.2, y + 0.1, 5.5, 0.5,
             size=15, bold=True, color=YELLOW)
    add_text(s8, desc, col + 0.2, y + 0.65, 5.5, 0.7, size=13, color=WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9: Erweiterungsmöglichkeiten
# ─────────────────────────────────────────────────────────────────────────────
s9 = prs.slides.add_slide(BLANK)
add_bg(s9)
add_rect(s9, 0, 0, 13.33, 1.1, BG_CARD)
add_rect(s9, 0, 1.08, 13.33, 0.05, ACCENT)
add_text(s9, "🚀  Nächste Schritte & Erweiterungen", 0.4, 0.2, 12, 0.7,
         size=28, bold=True, color=WHITE)

extensions = [
    ("☁️  Cloud-Deployment",    "Snowflake / BigQuery als Produktions-Data-Warehouse"),
    ("🤖  Machine Learning",    "Kursvorhersage-Modell auf Basis der berechneten Features"),
    ("📰  Sentiment-Analyse",   "Finanznachrichten als zusätzliche Datenquelle"),
    ("⚡  Echtzeit-Verarbeitung","Apache Kafka statt täglichem Batch-Betrieb"),
    ("🧪  Datentests",          "pytest + Great Expectations für robustere Qualitätsprüfung"),
    ("🔄  dbt Integration",     "SQL-Modelle als vollständige dbt-Pipeline"),
    ("🛸  Kubernetes",          "Container-Orchestrierung für Produktionsskalierung"),
    ("📊  BI-Integration",      "Metabase / Tableau direkt auf DuckDB"),
]
for i, (title, desc) in enumerate(extensions):
    col = 0.4 if i < 4 else 7.0
    row = i % 4
    y = 1.4 + row * 1.42
    add_text(s9, title, col, y, 5.8, 0.45, size=14, bold=True, color=ACCENT2)
    add_text(s9, f"   {desc}", col, y + 0.42, 5.8, 0.45, size=12, color=WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10: Abschluss
# ─────────────────────────────────────────────────────────────────────────────
s10 = prs.slides.add_slide(BLANK)
add_bg(s10)
add_rect(s10, 0, 0, 0.5, 7.5, ACCENT)
add_rect(s10, 0.5, 3.1, 12.83, 0.08, ACCENT)

add_text(s10, "🙏", 1.2, 0.8, 2, 1.2, size=60)
add_text(s10, "Vielen Dank!", 1.2, 1.8, 11, 1.0,
         size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s10, "Fragen & Diskussion", 1.2, 2.75, 11, 0.6,
         size=22, color=ACCENT, align=PP_ALIGN.LEFT)
add_text(s10, "GitHub: github.com/[username]/etl_pipeline", 1.2, 4.0, 11, 0.5,
         size=15, color=GRAY, align=PP_ALIGN.LEFT)
add_text(s10,
         "\"Dieses Projekt zeigt, dass ich den gesamten Data Engineering Stack beherrsche –\n"
         "von der Datenaufnahme über die Transformation bis zur Visualisierung.\"",
         1.2, 5.0, 11, 1.2, size=14, color=WHITE, align=PP_ALIGN.LEFT)

# ── Save ──────────────────────────────────────────────────────────────────────
OUTPUT = r"c:\etl_pipeline\Stock_ETL_Pipeline_Präsentation_v2.pptx"
prs.save(OUTPUT)
print(f"✅ PowerPoint gespeichert: {OUTPUT}")

import subprocess
subprocess.Popen(["start", OUTPUT], shell=True)
print("🚀 Datei wird geöffnet...")
